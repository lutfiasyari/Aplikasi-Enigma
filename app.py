from flask import Flask, request, render_template, send_file
from werkzeug.utils import secure_filename
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from fpdf import FPDF
from openpyxl import load_workbook, Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PyPDF2 import PdfWriter, PdfReader
from PIL import Image
import os
import base64
import io

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads/'
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'docx', 'pptx', 'xlsx', 'png', 'jpg', 'jpeg'}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Pastikan direktori uploads ada
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def vigenere_encrypt(plaintext, key):
    key = key.lower()
    encrypted = []
    key_length = len(key)
    for i, char in enumerate(plaintext):
        if char.isalpha():
            shift = ord(key[i % key_length]) - ord('a')
            if char.islower():
                encrypted.append(chr((ord(char) - ord('a') + shift) % 26 + ord('a')))
            else:
                encrypted.append(chr((ord(char) - ord('A') + shift) % 26 + ord('A')))
        else:
            encrypted.append(char)
    return ''.join(encrypted)

def vigenere_decrypt(ciphertext, key):
    key = key.lower()
    decrypted = []
    key_length = len(key)
    for i, char in enumerate(ciphertext):
        if char.isalpha():
            shift = ord(key[i % key_length]) - ord('a')
            if char.islower():
                decrypted.append(chr((ord(char) - ord('a') - shift) % 26 + ord('a')))
            else:
                decrypted.append(chr((ord(char) - ord('A') - shift) % 26 + ord('A')))
        else:
            decrypted.append(char)
    return ''.join(decrypted)

def rot13(text):
    return text.translate(str.maketrans("ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz",
                                          "NOPQRSTUVWXYZABCDEFGHIJKLMnopqrstuvwxyzabcdefghijklm"))

def caesar_cipher(text, shift):
    encrypted = []
    for char in text:
        if char.isalpha():
            shift_base = ord('a') if char.islower() else ord('A')
            encrypted.append(chr((ord(char) - shift_base + shift) % 26 + shift_base))
        else:
            encrypted.append(char)
    return ''.join(encrypted)

def write_to_file(content, filename, format_type):
    if format_type == 'txt':
        with open(filename, 'w') as f:
            f.write(content)
    elif format_type == 'pdf':
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, content, align="j")
        pdf.output(filename)
    elif format_type == 'docx':
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run(content)
        run.font.name = "Arial"
        run.font.size = Pt(12)
        para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY  # Rata tengah
        doc.save(filename)
    elif format_type in ['png', 'jpg', 'jpeg']:
        img_data = base64.b64decode(content)
        with open(filename, 'wb') as f:
            f.write(img_data)
    elif format_type == 'pptx':
        prs = Presentation()
    
    def split_text_efficiently(content, max_paragraphs_per_slide=26):
        # Memecah teks menjadi paragraf yang efisien
        paragraphs = [p.strip() for p in content.split('\n') if p.strip()]
        return paragraphs

    # Membagi konten
    content_paragraphs = split_text_efficiently(content)

    # Pengaturan slide dan text box
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)   # Lebar maksimal
    height = Inches(6)  # Tinggi maksimal

    current_slide = None
    current_text_frame = None

    for paragraph in content_paragraphs:
        # Buat slide baru jika belum ada atau text frame penuh
        if current_slide is None or (current_text_frame and len(current_text_frame.text) > 2000):
            current_slide = prs.slides.add_slide(prs.slide_layouts[5])
            content_box = current_slide.shapes.add_textbox(left, top, width, height)
            current_text_frame = content_box.text_frame
            current_text_frame.word_wrap = True

        # Tambahkan paragraf
        p = current_text_frame.add_paragraph()
        p.text = paragraph

        # Pengaturan font
        for run in p.runs:
            run.font.name = 'Times New Roman'
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0, 0, 0)
        
        p.alignment = PP_ALIGN.JUSTIFY

    prs.save(filename)

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        file = request.files['file']
        key = request.form['key']
        format_type = request.form['format_type']
        operation = request.form['operation']

        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)

            # Baca isi file
            if filename.endswith('.txt'):
                with open(file_path, 'r') as f:
                    content = f.read()
            elif filename.endswith('.pdf'):
                from PyPDF2 import PdfReader
                reader = PdfReader(file_path)
                content = '\n'.join(page.extract_text() for page in reader.pages if page.extract_text())
                content = ' '.join(content.split())
            elif filename.endswith('.docx'):
                doc = Document(file_path)
                content = '\n'.join([para.text for para in doc.paragraphs if para.text.strip()])
                content = ' '.join(content.split())
            elif filename.endswith('.png') or filename.endswith('.jpg') or filename.endswith('.jpeg'):
                with open(file_path, 'rb') as f:
                    img_data = f.read()
                    content = base64.b64encode(img_data).decode('utf-8')
            elif filename.endswith('.xlsx'):
                wb = load_workbook(file_path)
                content = ''
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        content += ' '.join(str(cell) for cell in row if cell is not None) + '\n'
            elif filename.endswith('.pptx'):
                prs = Presentation(file_path)
                content = ''
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                # Mengambil teks dan menghilangkan spasi berlebihan
                            text = shape.text.strip()
                            if text:  # Pastikan tidak menambahkan teks kosong
                                content += text + '\n'  # Menambahkan baris baru setelah setiap teks

    # Menghilangkan spasi berlebihan di akhir konten
                content = '\n'.join(line.strip() for line in content.splitlines() if line.strip())

            # Enkripsi atau dekripsi isi
            if operation == 'encrypt':
                if request.form['cipher'] == 'vigenere':
                    encrypted_content = vigenere_encrypt(content, key)
                elif request.form['cipher'] == 'rot13':
                    encrypted_content = rot13(content)
                elif request.form['cipher'] == 'caesar':
                    shift = int(key) 
                    encrypted_content = caesar_cipher(content, shift)
            elif operation == 'decrypt':
                if request.form['cipher'] == 'vigenere':
                    encrypted_content = vigenere_decrypt(content, key)
                elif request.form['cipher'] == 'rot13':
                    encrypted_content = rot13(content)
                elif request.form['cipher'] == 'caesar':
                    shift = int(key) 
                    encrypted_content = caesar_cipher(content, -shift)

            output_filename = os.path.join(app.config['UPLOAD_FOLDER'], 'output.' + format_type)
            write_to_file(encrypted_content, output_filename, format_type)

            return send_file(output_filename, as_attachment=True)

    return render_template('index.html')

if __name__ == '__main__':
    app.run(debug=True)